import uuid
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
import pdfplumber
import pdfkit
import imgkit
import datetime
from docx import Document
import pdfplumber
import re
import smtplib
import ssl
import pandas as pd
import string
import tkinter as tk
from tkinter import filedialog, StringVar, messagebox, ttk
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email.mime.text import MIMEText
from reportlab.pdfgen import canvas
from email.mime.image import MIMEImage
from email import encoders
import sys
from pptx import Presentation
from reportlab.lib.pagesizes import letter, landscape
from openpyxl import load_workbook
from docx2pdf import convert as docx2pdf_convert
from io import BytesIO
from PIL import Image
import customtkinter as ctk
import random
import subprocess
import urllib.request
import time
import os
import ctypes

def generate_unique_id(serial_number):
    prefix = f"I{serial_number:04d}"
    middle_part = ''.join(random.choices(string.ascii_uppercase + string.digits, k=7))
    suffix = ''.join(random.choices(string.ascii_uppercase + string.digits, k=12))
    unique_sequence = f"{prefix}_{middle_part}_{suffix}"
    return unique_sequence


def is_admin():
    """Check if the script is running with administrative privileges."""
    try:
        return ctypes.windll.shell32.IsUserAnAdmin()
    except:
        return False

def run_as_admin():
    """Run the script as an administrator."""
    ctypes.windll.shell32.ShellExecuteW(None, "runas", sys.executable, " ".join(sys.argv), None, 1)

def install_wkhtmltopdf(install_dir):
    """Install wkhtmltopdf and wkhtmltoimage."""
    print("wkhtmltopdf not found. Installing...")
    os.makedirs(install_dir, exist_ok=True)
    download_url = "https://github.com/wkhtmltopdf/packaging/releases/download/0.12.6-1/wkhtmltox-0.12.6-1.msvc2015-win64.exe"
    installer_path = os.path.join(install_dir, "wkhtmltox_installer.exe")

    print("Downloading wkhtmltopdf...")
    urllib.request.urlretrieve(download_url, installer_path)

    try:
        print("Installing wkhtmltopdf...")
        subprocess.run([installer_path, '/SILENT', '/VERYSILENT', f'/DIR={install_dir}'], check=True)
    finally:
        os.remove(installer_path)

# def get_executable_path(filename):
#     if getattr(sys, 'frozen', False):
#         # The application is frozen
#         base_path = sys._MEIPASS  # PyInstaller sets _MEIPASS to the temp directory where the bundled files are unpacked
#     else:
#         base_path = os.path.dirname(os.path.abspath(__file__)) 
#     executable_path = os.path.join(base_path, 'wkhtmltopdf', 'bin', filename) 
#     if not os.path.exists(executable_path):
#         raise FileNotFoundError(f"{filename} not found at {executable_path}. Ensure it is bundled correctly.")
#     return executable_path

def check_and_install_wkhtmltopdf():
    """Check if wkhtmltopdf and wkhtmltoimage are installed, otherwise install them."""
    base_path = r'C:\Program Files\wkhtmltopdf'
    wkhtmltopdf_path = os.path.join(base_path, 'bin', 'wkhtmltopdf.exe')
    wkhtmltoimage_path = os.path.join(base_path, 'bin', 'wkhtmltoimage.exe')

    if not os.path.isfile(wkhtmltopdf_path) or not os.path.isfile(wkhtmltoimage_path):
        install_wkhtmltopdf(base_path)

    return wkhtmltopdf_path, wkhtmltoimage_path

if not is_admin():
    run_as_admin()
    sys.exit()  # Exit script to rerun as admin

# Main setup
wkhtmltopdf_path, wkhtmltoimage_path = check_and_install_wkhtmltopdf()
config_pdfkit = pdfkit.configuration(wkhtmltopdf=wkhtmltopdf_path)
config_imgkit = imgkit.config(wkhtmltoimage=wkhtmltoimage_path)

print("wkhtmltopdf and wkhtmltoimage are set up successfully.")
def convert_txt_to_pdf(txt_file_path, pdf_file_path):
    c = canvas.Canvas(pdf_file_path, pagesize=letter)
    width, height = letter
    with open(txt_file_path, 'r') as f:
        text = f.readlines()
    y = height - 40
    for line in text:
        c.drawString(40, y, line.strip())
        y -= 15
    c.save()

def authenticate_gmail(credentials_path):
    if not isinstance(credentials_path, str):
        raise TypeError(f"Expected credentials_path to be a string, but got {type(credentials_path).__name__} instead.")
    creds = None
    SCOPES = ['https://www.googleapis.com/auth/gmail.send']  # Define the scope
    if os.path.exists('token.json'):
        creds = Credentials.from_authorized_user_file('token.json', SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            if not os.path.exists(credentials_path):
                raise FileNotFoundError(f"Credentials file not found: {credentials_path}")
            flow = InstalledAppFlow.from_client_secrets_file(credentials_path, SCOPES)
            creds = flow.run_local_server(port=0)
        with open('token.json', 'w') as token:
            token.write(creds.to_json())
    return creds

def convert_pptx_to_pdf(pptx_file_path, pdf_file_path):
    prs = Presentation(pptx_file_path)
    pdf_canvas = canvas.Canvas(pdf_file_path, pagesize=landscape(letter))
    for slide in prs.slides:
        pdf_canvas.showPage()
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image_stream = BytesIO(shape.image.blob)
                image = Image.open(image_stream)
                pdf_canvas.drawInlineImage(image, 40, 40)
    pdf_canvas.save()

def convert_xlsx_to_pdf(xlsx_file_path, pdf_file_path):
    workbook = load_workbook(xlsx_file_path)
    sheet = workbook.active
    pdf_canvas = canvas.Canvas(pdf_file_path, pagesize=landscape(letter))
    width, height = landscape(letter)
    y = height - 40
    for row in sheet.iter_rows(values_only=True):
        x = 40
        for cell in row:
            pdf_canvas.drawString(x, y, str(cell))
            x += 100
        y -= 20
        if y < 40:
            pdf_canvas.showPage()
            y = height - 40
    pdf_canvas.save()   

class BulkEmailApp:
    def __init__(self, master):
                self.root = root
                self.root.title("Bulk Email Sender")
                self.root.geometry("1300x600")

                # Initialize variables
                self.master = master
                self.sender_email_var = tk.StringVar()
                self.password_var = tk.StringVar()
                self.subject_var = tk.StringVar()
                self.text_size_var = tk.StringVar()
                self.recipient_var = tk.StringVar() 
                self.font_style_var = tk.StringVar()
                self.inline_image_var = tk.BooleanVar()
                self.convert_html_to_image_flag = tk.BooleanVar()
                self.send_as_html_var = tk.BooleanVar()
                self.convert_to_pdf_var = tk.BooleanVar()
                self.recipients = []
                self.files_path = []
                self.sender_emails = []
                self.sender_passwords = []
                self.sender_credentials = []  # Add this line
                self.credentials_path=""
                self.credentials_dir = ""
                self.credentials = []
                self.email_id = ""  # Initialize email ID
                self.Name = ""  # Initialize customer name
                self.company = ""  # Initialize company
                self.recipient = ""  # Initialize recipient

                self.create_widgets()

    def create_widgets(self):
        main_frame = ctk.CTkFrame(self.root)
        main_frame.grid(row=0, column=0, sticky="nsew")

        left_frame = ctk.CTkFrame(main_frame)
        left_frame.grid(row=0, column=0, padx=(10, 5), pady=10, sticky="nsew")

        right_frame = ctk.CTkFrame(main_frame)
        right_frame.grid(row=0, column=1, padx=(5, 10), pady=10, sticky="nsew")

        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        main_frame.columnconfigure(0, weight=1)
        main_frame.columnconfigure(1, weight=1)
        main_frame.rowconfigure(0, weight=1)
        main_frame.rowconfigure(1, weight=0)
        main_frame.rowconfigure(2, weight=1)
        left_frame.columnconfigure(0, weight=1)
        right_frame.columnconfigure(0, weight=1)

        # Upload Info block widgets
        ctk.CTkLabel(left_frame, text="Select File to Send:").grid(row=0, column=0, padx=5, pady=5, sticky='w')
        self.file_option = tk.StringVar(value="file")
        file_radio = ctk.CTkRadioButton(left_frame, text="File", variable=self.file_option, value="file", command=self.update_file_selection)
        file_radio.grid(row=0, column=1, padx=10, pady=5, sticky='w')
        directory_radio = ctk.CTkRadioButton(left_frame, text="Directory", variable=self.file_option, value="directory", command=self.update_file_selection)
        directory_radio.grid(row=0, column=2, padx=10, pady=5, sticky='w')

        select_button = ctk.CTkButton(left_frame, text="Browse", command=self.select_file)
        select_button.grid(row=0, column=3, padx=10, pady=5, sticky='w')
        
        ctk.CTkLabel(left_frame, text="Upload Email File:").grid(row=1, column=0, padx=5, pady=5, sticky='w')
        upload_file_button = ctk.CTkButton(left_frame, text="Browse", command=self.upload_file)
        upload_file_button.grid(row=1, column=3, padx=10, pady=5, sticky='w')

        ctk.CTkLabel(left_frame, text="Upload Credentials:").grid(row=2, column=0, padx=5, pady=5, sticky='w')
        # Add radio buttons for selecting credential upload type
        self.credential_option = tk.StringVar(value="file")
        credential_file_radio = ctk.CTkRadioButton(left_frame, text="File", variable=self.credential_option, value="file")
        credential_file_radio.grid(row=2, column=1, padx=10, pady=5, sticky='w')
        credential_dir_radio = ctk.CTkRadioButton(left_frame, text="Directory", variable=self.credential_option, value="directory")
        credential_dir_radio.grid(row=2, column=2, padx=10, pady=5, sticky='w')
        # Add upload credentials button
        ctk.CTkButton(left_frame, text="Browse", command=self.upload_credentials).grid(row=2, column=3, padx=10, pady=5, sticky='w')

        # Add a label and entry field for the sender's name
        ctk.CTkLabel(left_frame, text="Sender Name:").grid(row=3, column=0, padx=5, pady=5, sticky='w')
        self.sender_name_var = StringVar()
        sender_name_entry = ctk.CTkEntry(left_frame, textvariable=self.sender_name_var, width=200)
        sender_name_entry.grid(row=3, column=1, padx=10, pady=5, sticky='w')


        # Inline Image block widgets
        self.inline_image_var = tk.BooleanVar()
        inline_image_checkbutton = ctk.CTkCheckBox(left_frame, text="Inline Image", variable=self.inline_image_var)
        inline_image_checkbutton.grid(row=4, column=0, columnspan=2, padx=10, pady=5, sticky='w')

        # Convert HTML to Image checkbox
        convert_html_to_image_check = ctk.CTkCheckBox(left_frame, text="Convert HTML to Image", variable=self.convert_html_to_image_flag)
        convert_html_to_image_check.grid(row=5, column=0, columnspan=2, padx=10, pady=5, sticky="w")

        # Send as HTML checkbox
        self.send_as_html_var = tk.BooleanVar()
        send_as_html_check = ctk.CTkCheckBox(left_frame, text="Send as HTML", variable=self.send_as_html_var)
        send_as_html_check.grid(row=4, column=1, columnspan=2, padx=10, pady=5, sticky='w')

        self.convert_html_to_pdf_var = tk.BooleanVar()
        convert_html_to_pdf_checkbox = ctk.CTkCheckBox(left_frame, text="Convert HTML to PDF", variable=self.convert_html_to_pdf_var)
        convert_html_to_pdf_checkbox.grid(row=5, column=1, columnspan=2, padx=10, pady=5, sticky='w')

       # Sender Mail Field
        ctk.CTkLabel(right_frame, text="Upload Sender Mail File:").grid(row=0, column=0, padx=5, pady=5, sticky='w')
        upload_sender_mail_button = ctk.CTkButton(right_frame, text="Browse", command=self.upload_sender_mail_file)
        upload_sender_mail_button.grid(row=0, column=1, padx=10, pady=5, sticky='w')

        # Sender Password Field
        ctk.CTkLabel(right_frame, text="Upload Sender Password File:").grid(row=1, column=0, padx=5, pady=5, sticky='w')
        upload_sender_password_button = ctk.CTkButton(right_frame, text="Browse", command=self.upload_sender_password_file)
        upload_sender_password_button.grid(row=1, column=1, padx=10, pady=5, sticky='w')

        ctk.CTkLabel(right_frame, text="Email Subject:").grid(row=2, column=0, padx=10, pady=5, sticky='w')
        self.subject_var = StringVar()
        subject_entry = ctk.CTkEntry(right_frame, textvariable=self.subject_var, width=200)
        subject_entry.grid(row=2, column=1, padx=10, pady=5, sticky='w')

        ctk.CTkLabel(right_frame, text="Email Content:").grid(row=3, column=0, padx=10, pady=5, sticky='w,n')
        self.content_text = tk.Text(right_frame, width=40, height=10, wrap=tk.WORD)
        self.content_text.grid(row=3, column=1, padx=10, pady=5, sticky='w')

        # Options block widgets
        options_frame = ctk.CTkFrame(main_frame)
        options_frame.grid(row=1, column=0, columnspan=2, padx=10, pady=10, sticky="ew")

        # Convert to PDF checkbox
        self.convert_to_pdf_var = tk.BooleanVar()
        convert_to_pdf_check = ctk.CTkCheckBox(options_frame, text="Convert to PDF", variable=self.convert_to_pdf_var)
        convert_to_pdf_check.grid(row=0, column=0, padx=5, pady=5, sticky='w')

        ctk.CTkLabel(options_frame, text="Text Size:").grid(row=0, column=1, padx=10, pady=5, sticky='w')
        self.text_size_var = StringVar(value='12')  # Default text size
        text_size_dropdown = ctk.CTkOptionMenu(options_frame, variable=self.text_size_var, values=['8', '10', '12', '14', '16', '18', '20'])
        text_size_dropdown.grid(row=0, column=2, padx=10, pady=5, sticky='w')

        ctk.CTkLabel(options_frame, text="Font Style:").grid(row=0, column=3, padx=10, pady=5, sticky='w')
        self.font_style_var = StringVar(value='Arial')  # Default font style
        font_style_dropdown = ctk.CTkOptionMenu(options_frame, variable=self.font_style_var, values=['Arial', 'Times New Roman', 'Courier New'])
        font_style_dropdown.grid(row=0, column=4, padx=10, pady=5, sticky='w')

        # Create the log frame and TreeView with an additional 'Sender Email' column
        log_frame = ctk.CTkFrame(main_frame)
        log_frame.grid(row=2, column=0, columnspan=2, padx=10, pady=10, sticky="wn")

        columns = ('S.No', 'Email', 'Email ID', 'Name', 'Company', 'Sender Email', 'Status')
        self.log_treeview = ttk.Treeview(log_frame, columns=columns, show='headings')
        self.log_treeview.heading('S.No', text='S.No')
        self.log_treeview.heading('Email', text='Email')
        self.log_treeview.heading('Email ID', text='Email ID')
        self.log_treeview.heading('Name', text='Name')
        self.log_treeview.heading('Company', text='Company')
        self.log_treeview.heading('Sender Email', text='Sender Email')
        self.log_treeview.heading('Status', text='Status')

        self.log_treeview.grid(row=0, column=0, sticky='nsew')

        log_frame.columnconfigure(0, weight=1)
        log_frame.rowconfigure(0, weight=1)

        # Execution buttons
        buttons_frame = ctk.CTkFrame(main_frame)
        buttons_frame.grid(row=3, column=0, columnspan=2, padx=10, pady=10, sticky="ew")

        self.send_button = ctk.CTkButton(buttons_frame, text="Send Emails", command=self.send_emails)
        self.send_button.grid(row=0, column=0, padx=10, pady=5)

        preview_button = ctk.CTkButton(buttons_frame, text="Tags", command=self.preview_email)
        preview_button.grid(row=0, column=1, padx=10, pady=5)

        clear_button = ctk.CTkButton(buttons_frame, text="Clear", command=self.clear_fields)
        clear_button.grid(row=0, column=2, padx=10, pady=5)

        logout_button = ctk.CTkButton(buttons_frame, text="Close", command=self.logout)
        logout_button.grid(row=0, column=4, padx=10, pady=5)

        save_log_button = ctk.CTkButton(buttons_frame, text="Save Log to Excel", command=self.save_log_to_excel)
        save_log_button.grid(row=0, column=3, padx=10, pady=5)

        pass

    def upload_file(self):
        self.file_path = filedialog.askopenfilename(filetypes=[("Excel files", "*.xls *.xlsx")])
        if self.file_path:
            self.load_recipients(self.file_path)

    def load_recipients(self, file_path):
        try:
            df = pd.read_excel(file_path)
            self.recipients = df['Email'].tolist()
            self.recipient_details = {}  # Initialize the dictionary to store recipient details
            
            for _, row in df.iterrows():
                email = row['Email']
                name = row.get('Name')  # Fetch 'Name' column or default if not present
                company = row.get('Company')  # Fetch 'Company' column or default if not present
                address = row.get('Address')  # Fetch 'Address' column or default
                tfn_no = row.get('TFN No', '1-800-123-4567')  # Fetch 'TFN No' column or default
                
                # Store all details in the recipient_details dictionary
                self.recipient_details[email] = {
                    'name': name,
                    'company': company,
                    'address': address,
                    'tfn_no': tfn_no,
                }
            # Optionally show a success message
            # messagebox.showinfo("Success", f"Loaded {len(self.recipients)} recipients.")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to load recipients: {e}")
 
    def update_file_selection(self):
        pass
    
    def upload_sender_mail_file(self):
        file_path = filedialog.askopenfilename(filetypes=[("Excel files", "*.xlsx"), ("Text files", "*.txt")])
        if file_path:
            try:
                if file_path.endswith('.xlsx'):
                    df = pd.read_excel(file_path)
                    if 'Email' in df.columns:
                        self.sender_emails = df['Email'].tolist()
                    else:
                        messagebox.showerror("Error", "Excel file must contain an 'Email' column.")
                elif file_path.endswith('.txt'):
                    with open(file_path, 'r') as file:
                        self.sender_emails = [line.strip() for line in file]
                print("Sender emails loaded:", self.sender_emails)
                self.update_sender_credentials()  # Update credentials after loading emails
            except Exception as e:
                messagebox.showerror("Error", f"Failed to load sender emails: {e}")

    def upload_sender_password_file(self):
        file_path = filedialog.askopenfilename(filetypes=[("Excel files", "*.xlsx"), ("Text files", "*.txt")])
        if file_path:
            try:
                if file_path.endswith('.xlsx'):
                    df = pd.read_excel(file_path)
                    if 'Password' in df.columns:
                        self.sender_passwords = df['Password'].tolist()
                    else:
                        messagebox.showerror("Error", "Excel file must contain a 'Password' column.")
                elif file_path.endswith('.txt'):
                    with open(file_path, 'r') as file:
                        self.sender_passwords = [line.strip() for line in file]
                print("Sender passwords loaded:", self.sender_passwords)
                self.update_sender_credentials()  # Update credentials after loading passwords
            except Exception as e:
                messagebox.showerror("Error", f"Failed to load sender passwords: {e}")
    
    def update_sender_credentials(self):
        self.sender_credentials = list(zip(self.sender_emails, self.sender_passwords))
        print(type(self.sender_credentials))
        print(self.sender_credentials)

    def select_file(self):
        file_option = self.file_option.get()
        if file_option == "file":
            file_path = filedialog.askopenfilename()
            self.files_path = [file_path] if file_path else []
        elif file_option == "directory":
            directory_path = filedialog.askdirectory()
            if directory_path:
                self.files_path = [os.path.join(directory_path, f) for f in os.listdir(directory_path) if os.path.isfile(os.path.join(directory_path, f))]
            else:
                self.files_path = []
        
    def generate_email_id(self):
        return ''.join(random.choices(string.ascii_letters + string.digits, k=8))

    def upload_credentials(self):
        if self.credential_option.get() == "file":
            self.credentials_path = filedialog.askopenfilename(filetypes=[("JSON files", "*.json")])
            if self.credentials_path:
                print(f"Selected credentials file: {self.credentials_path}")  # Debugging line
                self.load_credentials([self.credentials_path])
        elif self.credential_option.get() == "directory":
            self.credentials_dir = filedialog.askdirectory()
            if self.credentials_dir:
                json_files = [os.path.join(self.credentials_dir, file) for file in os.listdir(self.credentials_dir) if file.endswith('.json')]
                print(f"Selected credentials directory: {self.credentials_dir}")  # Debugging line
                print(f"JSON files in directory: {json_files}")  # Debugging line
                self.load_credentials(json_files)

    def load_credentials(self, file_paths):
        self.credentials = []
        for file_path in file_paths:
            if os.path.exists(file_path):
                self.credentials.append(file_path)  # Store file paths instead of JSON content
                print(f"Loaded credentials from {file_path}")  # Debugging line
            else:
                messagebox.showerror("Error", f"Credentials file not found: {file_path}")

    def authenticate_and_initialize_service(self):
        if not self.credentials:
            messagebox.showerror("Error", "No credentials loaded. Please upload credentials first.")
            return 
        if not self.credentials_path and not self.credentials_dir:
            messagebox.showerror("Error", "No credentials path specified.")
            return

        self.sender_services = []
        for credential_path in self.credentials:
            print(f"Authenticating with credential path: {credential_path}")
            creds = authenticate_gmail(credential_path)  # Pass the file path here
            if creds:
                print(f"Authenticated successfully with {credential_path}")
                self.sender_services.append(creds)
            else:
                print(f"Authentication failed for {credential_path}")
    
    def convert_html_to_pdf(self, html_file, pdf_path):
        # Verbose output for debugging
        try:
            # Generate the PDF
            pdfkit.from_file(html_file, pdf_path, configuration=config_pdfkit, verbose=True)
            print(f"PDF generated successfully at {pdf_path}")
        except OSError as e:
            print(f"Error generating PDF: {e}")

        # Optional: Generate the CLI command for debugging
        try:
            pdf_kit_instance = pdfkit.PDFKit(html_file, 'file', configuration=config_pdfkit)
            cli_command = ' '.join(pdf_kit_instance.command())
            print(f"Generated wkhtmltopdf command: {cli_command}")
        except OSError as e:
            print(f"Error generating wkhtmltopdf command: {e}")

    def convert_html_url_to_image(self, html_file_path, image_file_path):
        try:
            # Read and print the HTML content for debugging
            with open(html_file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
                print("HTML Content:", html_content)  # Debugging line to check HTML content
            # Convert HTML to image
            imgkit.from_file(html_file_path, image_file_path, config=config_imgkit)
            print(f"Image generated successfully at {image_file_path}")
        except OSError as e:
            print(f"Error generating image: {e}")
    
    
    def send_emails(self):
        subject = self.subject_var.get()
        content = self.content_text.get("1.0", tk.END).strip()
        text_size = self.text_size_var.get()
        font_style = self.font_style_var.get()
        inline_image = self.inline_image_var.get()
        convert_html_to_image_flag = self.convert_html_to_image_flag.get()
        send_as_html = self.send_as_html_var.get()
        convert_to_pdf = self.convert_to_pdf_var.get()
        convert_html_to_pdf = self.convert_html_to_pdf_var.get()
        sender_name = self.sender_name_var.get()  # Get the sender's name

        if not self.recipients or not subject or not content:
            messagebox.showerror("Error", "Please fill in all fields and upload the necessary files.")
            return

        num_senders = len(self.sender_credentials)
        num_recipients = len(self.recipients)
        recipients_per_sender = num_recipients // num_senders
        extra_recipients = num_recipients % num_senders

        recipient_splits = []
        start_index = 0
        for i in range(num_senders):
            end_index = start_index + recipients_per_sender + (1 if i < extra_recipients else 0)
            recipient_splits.append(self.recipients[start_index:end_index])
            start_index = end_index

        port = 465
        smtp_server = "smtp.gmail.com"
        context = ssl.create_default_context()
        all_sent_successfully = True
        self.authenticate_and_initialize_service()
        serial_number = 1
        sender_index = 0

        for creds, sender_credential in zip(self.sender_services, self.sender_credentials):
            sender_email, sender_password = sender_credential
            if not sender_password:
                continue
            recipient_list = recipient_splits[sender_index]
            sender_recipients = recipient_splits[sender_index]

            for recipient in sender_recipients:
                recipient_details = self.recipient_details.get(recipient, {})
                name = recipient_details.get('name', '')
                company = recipient_details.get('company', '')
                email_id = self.generate_email_id()

                # Generate a unique identifier for this email
                unique_id = generate_unique_id(serial_number)

                # Create the email message
                message = MIMEMultipart("related")
                message["Subject"] = subject
                if sender_name:
                    message["From"] = f"{sender_name} <{sender_email}>"
                else:
                    message["From"] = sender_email
                message["To"] = recipient
                message["Disposition-Notification-To"] = sender_email
                attached_files = set()

                combined_content = ""
                html_content = ""
                # Load the HTML content if 'Send as HTML' is checked
                if send_as_html:
                    for file in self.files_path:
                        if file.lower().endswith('.html'):
                            with open(file, 'r', encoding='utf-8') as html_file:
                                html_content = html_file.read()
                            attached_files.add(file)
                            break

                # Construct the combined content
                if html_content:
                    combined_content = f"""
                    <p>{content}</p>

                    
                    <html>
                    <body style="font-size:{text_size}px; font-family:{font_style}; margin: 0; padding: 0;">
                    {unique_id}
                        <div style="max-width: 800px; margin: auto; padding: 20px; border: 1px solid #ddd; border-radius: 10px; background-color: #f9f9f9;">
                            {html_content}
                            """
                    if inline_image:
                        combined_content += '<img src="cid:inline_image" style="max-width: 100%; height: auto;"><br>'
                    combined_content += """
                        </div>
                    </body>
                    </html>
                    """
                else:
                    combined_content = f"""
                    <html>
                    <body style="font-size:{text_size}px; font-family:{font_style};">
                        <p>{content}</p>
                        """
                    if inline_image:
                        combined_content += '<img src="cid:inline_image" style="max-width: 100%; height: auto;"><br>'
                    combined_content += """
                    </body>
                    </html>
                    """
                
                message.attach(MIMEText(combined_content, 'html'))

                # Attach inline image if required
                if inline_image:
                    for file in self.files_path:
                        if file.lower().endswith(('jpg', 'jpeg', 'png')):
                            with open(file, 'rb') as img:
                                image = MIMEImage(img.read())
                                image.add_header('Content-ID', '<inline_image>')
                                image.add_header('Content-Disposition', 'inline', filename=os.path.basename(file))
                                message.attach(image)
                            attached_files.add(file)

                if convert_html_to_image_flag:
                    for file in self.files_path:
                        if file.lower().endswith('.html') and file not in attached_files:
                            image_path = os.path.splitext(file)[0] + f'_{unique_id}.png'
                            self.convert_html_url_to_image(file, image_path)
                            if os.path.exists(image_path):
                                with open(image_path, 'rb') as img_file:
                                    image_content = img_file.read()
                                image = MIMEImage(image_content)
                                content_id = 'html_image'
                                image.add_header('Content-ID', f'<{content_id}>')
                                image.add_header('Content-Disposition', 'inline', filename=os.path.basename(image_path))
                                message.attach(image)
                                attached_files.add(file)  
                            else:
                                print(f"Failed to create image from {file}.")
                                messagebox.showerror("Error", f"Failed to create image from {file}.")
                                return
                if convert_html_to_pdf:
                    for file in self.files_path:
                        if file.lower().endswith('.html') and file not in attached_files:
                            pdf_path = os.path.splitext(file)[0] + f'_{unique_id}.pdf'
                            self.convert_html_to_pdf(file, pdf_path)
                            
                            with open(pdf_path, "rb") as pdf_file:
                                part = MIMEBase("application", "octet-stream")
                                part.set_payload(pdf_file.read())
                            encoders.encode_base64(part)
                            part.add_header('Content-Disposition', f'attachment; filename={os.path.basename(pdf_path)}')
                            message.attach(part)
                            attached_files.add(file)
                            break
                for file in self.files_path:
                    if file in attached_files:
                        continue
                    new_file_name = f"{os.path.splitext(os.path.basename(file))[0]}_{unique_id}{os.path.splitext(file)[1]}"
                    if convert_to_pdf:
                        if file.lower().endswith(('jpg', 'jpeg', 'png')):
                            image = Image.open(file)
                            pdf_bytes = BytesIO()
                            image.save(pdf_bytes, format='PDF')
                            pdf_bytes.seek(0)
                            part = MIMEBase('application', 'octet-stream')
                            part.set_payload(pdf_bytes.read())
                            encoders.encode_base64(part)
                            part.add_header('Content-Disposition', f'attachment; filename={new_file_name}.pdf')
                            message.attach(part)
                            attached_files.add(file)
                        elif file.lower().endswith('.txt'):
                            pdf_path = os.path.splitext(file)[0] + f'_{unique_id}.pdf'
                            convert_txt_to_pdf(file, pdf_path)
                            with open(pdf_path, "rb") as f:
                                part = MIMEBase("application", "octet-stream")
                                part.set_payload(f.read())
                            encoders.encode_base64(part)
                            part.add_header('Content-Disposition', f'attachment; filename={os.path.basename(pdf_path)}')
                            message.attach(part)
                            attached_files.add(file)
                        elif file.lower().endswith('.docx'):
                            pdf_path = os.path.splitext(file)[0] + f'_{unique_id}.pdf'
                            docx2pdf_convert(file, pdf_path)
                            with open(pdf_path, "rb") as f:
                                part = MIMEBase("application", "octet-stream")
                                part.set_payload(f.read())
                            encoders.encode_base64(part)
                            part.add_header('Content-Disposition', f'attachment; filename={os.path.basename(pdf_path)}')
                            message.attach(part)
                            attached_files.add(file)
                        elif file.lower().endswith('.pptx'):
                            pdf_path = os.path.splitext(file)[0] + f'_{unique_id}.pdf'
                            convert_pptx_to_pdf(file, pdf_path)
                            with open(pdf_path, "rb") as f:
                                part = MIMEBase("application", "octet-stream")
                                part.set_payload(f.read())
                            encoders.encode_base64(part)
                            part.add_header('Content-Disposition', f'attachment; filename={os.path.basename(pdf_path)}')
                            message.attach(part)
                            attached_files.add(file)
                        elif file.lower().endswith('.xlsx'):
                            pdf_path = os.path.splitext(file)[0] + f'_{unique_id}.pdf'
                            convert_xlsx_to_pdf(file, pdf_path)
                            with open(pdf_path, "rb") as f:
                                part = MIMEBase("application", "octet-stream")
                                part.set_payload(f.read())
                            encoders.encode_base64(part)
                            part.add_header('Content-Disposition', f'attachment; filename={os.path.basename(pdf_path)}')
                            message.attach(part)
                            attached_files.add(file)
                    else:
                        with open(file, "rb") as f:
                            part = MIMEBase("application", "octet-stream")
                            part.set_payload(f.read())
                        encoders.encode_base64(part)
                        part.add_header('Content-Disposition', f'attachment; filename={new_file_name}')
                        message.attach(part)
                        attached_files.add(file)

                try:
                    with smtplib.SMTP_SSL(smtp_server, port, context=context) as server:
                        server.login(sender_email, sender_password)
                        server.sendmail(sender_email, recipient, message.as_string())
                    self.log_treeview.insert('', 'end', values=(serial_number, recipient, email_id, name, company, sender_email, 'Sent'))
                    serial_number += 1
                except Exception as e:
                    all_sent_successfully = False
                    self.log_treeview.insert('', 'end', values=(serial_number, recipient, email_id, name, company, sender_email, f'Failed: {e}'))
                    serial_number += 1
            sender_index += 1
        if all_sent_successfully:
            messagebox.showinfo("Success", "Emails sent successfully.")
        else:
            messagebox.showerror("Error", "Some emails failed to send.")

    def preview_email(self):
        if not self.recipients:
            messagebox.showwarning("Warning", "No recipients available to preview.")
            return

        recipient = self.recipients[0]  # Select the first recipient or modify as per your logic
        recipient_details = self.recipient_details.get(recipient, {})
        
        # Fetching dynamic details
        name = recipient_details.get('name',)
        company = recipient_details.get('company')
        email_id = recipient  # Assuming recipient is the email ID
        address = recipient_details.get('address')  # Default value
        tfn_no = recipient_details.get('tfn_no', '1-800-123-4567')  # Default value
        sender_name = self.sender_credentials[0][0] 
        
        # Generating dynamic values
        invoice_number = self.generate_random_invoice_number()
        tran_number = self.generate_random_transaction_number()
        amount_no = self.generate_random_amount()
        cx_date = self.get_today_date()
        key_rnd = str(uuid.uuid4())
        cus_id = self.generate_random_customer_id()
        random_value = self.generate_random_alphanumeric()
# Company:         {company}
#         Address:         {address} ,Email ID:        {email_id}
        preview_text = f"""

        Customer Name:  {name}
        Customer Email: {email_id}
        Sender email:   {sender_name}
        TFN No:         {tfn_no}
        Address:        {address}
        """
        preview_text += f"""
        Invoice Number:  {invoice_number}
        Transaction Number: {tran_number}
        Amount:          {amount_no}
        Date:            {cx_date}
        Random Key:      {key_rnd}
        Customer ID:     {cus_id}
        Random Value:    {random_value}
        """

        if self.files_path:
            for file in self.files_path:
                if file.lower().endswith(('.html', '.txt', '.docx', '.pdf')):  # Extend as needed
                    content = self.get_file_content(file)
                    extracted_content =self.extract_relevant_content_with_replaced_tags(content, {
                    # "$cx_email": email_id,
                    # "$cx_name": name,
                    "Address": recipient_details.get('address', 'Unknown Address'),
                    "Company": company,
                    # "$tfn_no": recipient_details.get('tfn_no', 'Unknown TFN'),
                    # "$Sender_name": "Your Sender Name",
                    # "$invoice_number": invoice_number,
                    # "$amount_no": amount_no,
                    # "Date": cx_date,
                    "Tran_Number": tran_number,
                    # "Key_rnd": key_rnd,
                    # "$cus_id": cus_id,
                    # "$random": random_value,
                })
                    preview_text += f"""
                Extracted Content with Replaced Tags:
                {extracted_content}
                """
                break  # Stop after processing the first valid file

        messagebox.showinfo("Email Preview", preview_text.strip())

    def extract_relevant_content_with_replaced_tags(self, content, tag_map):
        """Replaces the tags in the content and extracts relevant sections."""
        relevant_content = ""
        
        # Replace tags with their corresponding values
        for tag, value in tag_map.items():
            if tag in content:
                relevant_content += f"{tag} : {value}\n"
                content = content.replace(tag, value)
        return relevant_content or "No tags were found in the content."
    
    def generate_random_invoice_number(self):
        return f"{random.randint(10, 99)}{self.generate_random_alphanumeric(6)}{random.randint(100, 999)}"

    def generate_random_transaction_number(self):
        return str(random.randint(100000000000, 999999999999))

    def generate_random_amount(self):
        return str(random.randint(199, 999))

    def get_today_date(self):
        return datetime.datetime.now().strftime('%d %b-%Y')

    def generate_random_customer_id(self):
        return f"CX{random.randint(1000000, 9999999)}"

    def generate_random_alphanumeric(self, length=7):
        chars = "ABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789"
        return ''.join(random.choices(chars, k=length))

    def extract_detail_from_file(self, file_content, detail_label):
        pattern = fr"{detail_label}\s*[:\s]*([\w/\-\s]+)"
        match = re.search(pattern, file_content, re.IGNORECASE)
        return match.group(1).strip() if match else None

    def get_file_content(self, file_path):
        content = ""
        if file_path.lower().endswith('.html'):
            content = self.get_html_content(file_path)
        elif file_path.lower().endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
        elif file_path.lower().endswith('.docx'):
            doc = Document(file_path)
            content = '\n'.join([para.text for para in doc.paragraphs])
        elif file_path.lower().endswith('.pdf'):
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    content += page.extract_text()
        return content

    def get_html_content(self, file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            print(f"Error reading HTML file: {e}")
            return ""  

    def clear_fields(self):
        self.subject_var.set("")  # Clear the subject entry field
        self.content_text.delete(1.0, tk.END)  # Clear the content text field
        self.recipient_var.set("")  # Clear the recipient entry field
        self.text_size_var.set("12")  # Reset to default text size
        self.font_style_var.set("Arial")  # Reset to default font style
        self.inline_image_var.set(False)  # Uncheck inline image option
        self.convert_html_to_image_flag.set(False)  # Uncheck HTML to image conversion option
        self.send_as_html_var.set(False)  # Uncheck send as HTML option
        self.convert_to_pdf_var.set(False)  # Uncheck convert to PDF option
        self.files_path = []
        self.recipients = []  # Assuming self.recipients is the list used
        self.sender_credentials = []  # Clear the sender credentials
        self.email_id = ""  # Clear any generated email ID
        self.Name = ""  # Clear customer name
        self.company = ""  # Clear company if used
        self.recipient = ""  # Clear recipient

    def logout(self):
        self.root.quit() 

    def save_log_to_excel(self):
        log_file = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel files", "*.xlsx")])
        if not log_file:
            return
        columns = [col for col in self.log_treeview["columns"]]
        log_data = [self.log_treeview.item(item)["values"] for item in self.log_treeview.get_children()]
        df = pd.DataFrame(log_data, columns=columns)
        df.to_excel(log_file, index=False)
        messagebox.showinfo("Success", "Log saved to Excel successfully.")

if __name__ == "__main__":
    ctk.set_appearance_mode("dark")  # Set the application to dark mode
    root = ctk.CTk()
    app = BulkEmailApp(root)
    root.mainloop()
